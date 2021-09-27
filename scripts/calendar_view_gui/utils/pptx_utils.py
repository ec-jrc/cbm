#!/usr/bin/env python3
# -*- coding: utf-8 -*-

# This file is part of CbM (https://github.com/ec-jrc/cbm).
# Author    : Daniele Borio
# Credits   : GTCAP Team
# Copyright : 2021 European Commission, Joint Research Centre
# License   : 3-Clause BSD
# Version   : Created on Mon Sep  7 10:00:19 2020


# Import python library for power point manipulation
import pptx
import matplotlib.pyplot as plt
import os
import batch_utils

"""
    Summary: 
        Functon that add a slides to the a power-point presentation using
        the NDVI profile and the imaginette summary generated from Sentinel data
        
    Arguments :
        prs - presentation object (from python-pptx)
        pptx_filename - where to save the presentation
        out_tif_folder_base - pathname of the output tif (string)
        parcel_id - parcel ID as a string
        crop - declared crop type provided as a string
        buffer_size_meter - buffer size in metres (int)
        vector_color - vector color as string - these arguments are used to build the 
                       pathnames of the JPG to be included in the slide
"""

def add_slide( prs, out_tif_folder_base, parcel_id, crop, buffer_size_meter, vector_color ) :
    
    # Open the presentation to be modified
    # prs = pptx.Presentation(pptx_filename)
    
    # build the address of the NDVI JPG
    path = out_tif_folder_base + '/ndvi_graphs/'
    
    # # sorry for patching it here (by Csaba), we could add it as a parameter later
    # path = out_tif_folder_base + '/ndvi_graphs_with_mean_new/'
    
    # convert parcel_id to a string that can be used as filename           
    parcel_id_as_filename = batch_utils.convert_string_to_filename(parcel_id)
    
    if parcel_id.isnumeric() :
        parcel_id_as_num = int(parcel_id)   # this is to remove initial 0s
        parcel_id_as_str = str(parcel_id_as_num)
    else :
        parcel_id_as_str = parcel_id
    
    # search for file 
    ndvi_files = [f for f in os.listdir(path) if parcel_id_as_filename in f and f.endswith('_NDVI.jpg') ]
    
    if len(ndvi_files) == 0 :
        return
    else :
        ndvi_jpg = path + ndvi_files[ 0 ] 
        
        
    
    # build the address of the imaginettes summary 
    summary_jpg = out_tif_folder_base + '/overview_jpg_half_weekly/'
    summary_jpg += ( parcel_id_as_filename + '_')
    summary_jpg += ( crop + '_')
    summary_jpg += "buff%dm_" % buffer_size_meter
    summary_jpg += vector_color + '.jpg'
    
    ############## Now build the new slide ###################################
    blank_slide_layout = prs.slide_layouts[14]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    
    # NDVI IMAGE
    ndvi_height = pptx.util.Cm(7.6)
    ndvi_width = pptx.util.Cm(13.05)
    ndvi_top = pptx.util.Cm(3.5)
    ndvi_left = pptx.util.Cm(1)
    
    slide.shapes.add_picture(ndvi_jpg, ndvi_left, ndvi_top, height=ndvi_height, width= ndvi_width)
    
    # Imaginettes Summary
    
    # Open the image to check if the width is within the slide
    cview = plt.imread( summary_jpg )
    
    height = cview.shape[ 0 ]
    width = cview.shape[ 1 ]
    
    # 19.0 is the new height in the slide
    scaled_width = float( width ) * 19.0 / height
    
    summary_left = pptx.util.Cm(14)
    
    # 19.9 is the maximum width allowed
    if scaled_width < 19.9 :
        summary_height = pptx.util.Cm(19.0)
        summary_top = pptx.util.Cm(0)
    else :
        scaled_height = ( 19.9 / width ) * height
        summary_height = pptx.util.Cm( scaled_height )
        # Center the image
        summary_top = pptx.util.Cm( (19.0 - scaled_height) / 2 )
    
    slide.shapes.add_picture(summary_jpg, summary_left, summary_top, height=summary_height)
    
    # Add textbox with parcel crop/information
    txt_left = pptx.util.Cm(0.5)
    txt_top = pptx.util.Cm(0.5)
    txt_height = pptx.util.Cm(1.8)
    txt_width = pptx.util.Cm(14.11)
    
    txBox = slide.shapes.add_textbox(txt_left, txt_top, txt_width, txt_height)
    tf = txBox.text_frame
    
    tf.text = "Parcel id = " + parcel_id
    p = tf.add_paragraph()
    p.text = "Declared crop = " + crop
    